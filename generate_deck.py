import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # White, Black & Blue Color Palette (Clean, Crisp, Professional)
    # -------------------------------------------------------------
    COLOR_BG = RGBColor(255, 255, 255)            # Pure White #FFFFFF
    COLOR_CARD_BG = RGBColor(248, 250, 252)       # Light Slate #F8FAFC
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)   # Border #E2E8F0
    COLOR_PRIMARY_BLUE = RGBColor(37, 99, 235)    # Royal Blue #2563EB
    COLOR_DARK_BLUE = RGBColor(30, 58, 138)       # Deep Navy #1E3A8A
    COLOR_BLACK = RGBColor(15, 23, 42)            # Charcoal Black #0F172A
    COLOR_MUTED = RGBColor(100, 116, 139)         # Slate 500 #64748B
    COLOR_ROSE = RGBColor(225, 29, 72)            # Crimson #E11D48
    COLOR_EMERALD = RGBColor(5, 150, 105)         # Forest Emerald #059669
    COLOR_LIGHT_BLUE = RGBColor(239, 246, 255)    # Soft Blue Tint #EFF6FF
    COLOR_BLUE_BORDER = RGBColor(191, 219, 254)   # Blue Border #BFDBFE

    def add_header(slide, tag_text, title_text, subtitle_text=""):
        # White Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Top Accent Blue Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
        bar.line.fill.background()

        # Tag / Category Pill
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.4))
        tf_tag = tb_tag.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_PRIMARY_BLUE
        p_tag.font.name = "Calibri"

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BLACK
        p_title.font.name = "Calibri"

        # Subtitle
        if subtitle_text:
            p_sub = tf_title.add_paragraph()
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = COLOR_MUTED
            p_sub.font.name = "Calibri"

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Hero)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()

    # Top & Bottom Blue Accent Ribbons
    t_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    t_bar.fill.solid()
    t_bar.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
    t_bar.line.fill.background()

    # Hero Card Container
    hero_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.8), Inches(10.933), Inches(5.9))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    hero_card.line.color.rgb = COLOR_BLUE_BORDER

    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.1), Inches(10.333), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p00 = tf.paragraphs[0]
    p00.text = "ZERO-TRUST RECRUITMENT DEFENSE"
    p00.font.size = Pt(12)
    p00.font.bold = True
    p00.font.color.rgb = COLOR_PRIMARY_BLUE
    p00.font.name = "Calibri"
    p00.alignment = PP_ALIGN.CENTER

    p0 = tf.add_paragraph()
    p0.text = "🛡️ HIRESHIELD"
    p0.font.size = Pt(46)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_DARK_BLUE
    p0.font.name = "Calibri"
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "Autonomous AI Job Scam Intelligence & Threat Verification Platform"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_BLACK
    p1.font.name = "Calibri"
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "\nDeterministic Risk Scoring • Multi-Modal OCR • Real-Time DNS/MX Telemetry • Cryptographic Job Passports"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_MUTED
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\n🌐 Live Deployed App: https://hireshield-ro6i.onrender.com"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_PRIMARY_BLUE
    p3.font.name = "Calibri"
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "🎥 Demo Video (Google Drive): https://drive.google.com/file/d/1MCPUq-A6OGeyXfKM-25laUY6rRLW4mTh/view?usp=drive_link"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_DARK_BLUE
    p4.font.name = "Calibri"
    p4.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 2: The Problem (Threat Landscape)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "The Threat Landscape", "The $3.8B Recruitment Fraud Epidemic", "Job seekers and early-career candidates face organized, highly lucrative cybercrime.")

    cards2 = [
        ("💸 Advance-Fee & Equipment Check Scams", "Scammers send counterfeit $3,000+ PDF checks, instructing victims to deposit and wire $400 via Zelle/Wire to a fake hardware vendor before the check bounces.", COLOR_ROSE),
        ("🌐 Typo-Squatted & Disposable Domains", "Attackers register lookalike domains on high-risk TLDs (.top, .xyz, .click) such as stripe-careers.top to host fraudulent application forms.", COLOR_PRIMARY_BLUE),
        ("💬 Off-Platform Telegram & Chat Traps", "Recruiters solicit sensitive candidate PII (SSNs, banking details) via unmonitored Telegram/WhatsApp channels to evade enterprise audit logs.", COLOR_DARK_BLUE)
    ]
    for i, (title, desc, color) in enumerate(cards2):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.1)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        ctb = slide2.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), Inches(3.2), Inches(3.9))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        cp0 = ctf.paragraphs[0]
        cp0.text = title
        cp0.font.size = Pt(16)
        cp0.font.bold = True
        cp0.font.color.rgb = color
        cp0.font.name = "Calibri"

        cp1 = ctf.add_paragraph()
        cp1.text = f"\n{desc}"
        cp1.font.size = Pt(13)
        cp1.font.color.rgb = COLOR_BLACK
        cp1.font.name = "Calibri"

    # -------------------------------------------------------------
    # SLIDE 3: The Solution (Architecture)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "System Architecture", "The HireShield 5-Stage Multi-Vector Pipeline", "Deterministic verification replacing black-box AI guesswork.")

    steps = [
        ("Stage 1 • Multi-Modal Ingestion", "Raw job text, job URLs, offer letter PDFs, and screenshots (Pytesseract OCR + Gemini Vision)."),
        ("Stage 2 • Dual-Engine Extraction", "Google Gemini 1.5 structured JSON entity parsing with zero-fail deterministic regex fallback."),
        ("Stage 3 • Deep DNS & MX Forensics", "Real-time corporate DNS MX record lookups, domain-to-email mismatch detection, and TLD reputation checks."),
        ("Stage 4 • Deterministic Risk Engine", "Explainable 0-100 mathematical risk deduction ledger (-40, -25, -20, -15 penalties with evidence trails)."),
        ("Stage 5 • Cryptographic Job Passport™", "Tamper-evident SHA-256 digital certificate generation with verifiable security seals and PDF export.")
    ]
    for i, (step_title, step_desc) in enumerate(steps):
        top = Inches(1.95 + i * 0.98)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        stb = slide3.shapes.add_textbox(Inches(1.0), top + Inches(0.12), Inches(11.3), Inches(0.65))
        stf = stb.text_frame
        stf.word_wrap = True
        sp0 = stf.paragraphs[0]
        sp0.text = f"{step_title}: "
        sp0.font.bold = True
        sp0.font.size = Pt(14)
        sp0.font.color.rgb = COLOR_PRIMARY_BLUE
        sp0.font.name = "Calibri"

        run = sp0.add_run()
        run.text = step_desc
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = COLOR_BLACK

    # -------------------------------------------------------------
    # SLIDE 4: Deterministic Scoring vs Black-Box AI
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Core Innovation", "Why Deterministic Scoring Beats Black-Box AI", "Transparent mathematical deduction builds uncompromising user trust.")

    # Left Column: Black Box AI
    left_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(254, 242, 242)
    left_card.line.color.rgb = RGBColor(254, 202, 202)

    ltb = slide4.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5.0), Inches(4.1))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    lp0 = ltf.paragraphs[0]
    lp0.text = "❌ Typical Black-Box LLM Approach"
    lp0.font.size = Pt(17)
    lp0.font.bold = True
    lp0.font.color.rgb = COLOR_ROSE

    lp1 = ltf.add_paragraph()
    lp1.text = "\n• Vague warnings: 'This email feels suspicious.'\n• Hallucinations & inconsistent scores on identical inputs.\n• No verifiable network or MX domain evidence.\n• Total system failure when API quotas or network drop."
    lp1.font.size = Pt(13)
    lp1.font.color.rgb = COLOR_BLACK

    # Right Column: HireShield Deterministic
    right_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.1), Inches(5.7), Inches(4.6))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    right_card.line.color.rgb = COLOR_BLUE_BORDER

    rtb = slide4.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.1), Inches(4.1))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    rp0 = rtf.paragraphs[0]
    rp0.text = "✅ HireShield Zero-Trust Engine"
    rp0.font.size = Pt(17)
    rp0.font.bold = True
    rp0.font.color.rgb = COLOR_DARK_BLUE

    rp1 = rtf.add_paragraph()
    rp1.text = "\n• Base 100 PTS with explicit mathematical deductions:\n  - -40 PTS: Advance Equipment Check Demand\n  - -25 PTS: High-Risk .top/.xyz Domain\n  - -20 PTS: Recruiter MX vs Corporate Domain Mismatch\n• Live DNS MX server handshake proof.\n• Zero-fail client & server offline resilience."
    rp1.font.size = Pt(13)
    rp1.font.color.rgb = COLOR_BLACK

    # -------------------------------------------------------------
    # SLIDE 5: Key Live Features & UX
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Product Capabilities", "Comprehensive Feature Suite & User Experience", "Engineered for intuitive candidate safety across web and mobile viewports.")

    feat_cards = [
        ("🖥️ 3D Interactive Cyber Terminal", "Cinematic CSS 3D clamshell laptop viewport with gyroscopic mouse-tracking telemetry.", COLOR_PRIMARY_BLUE),
        ("📸 Multi-Modal OCR Scanner", "Drop screenshots of offer letters or recruiter chat transcripts for instantaneous entity extraction.", COLOR_DARK_BLUE),
        ("🏢 Verified Enterprise Directory", "Search pre-verified companies (Stripe, Google, Microsoft, Razorpay) with active MX records.", COLOR_PRIMARY_BLUE),
        ("📜 Job Trust Passport™", "Export tamper-proof PDF digital credentials signed with unique SHA-256 cryptographic hashes.", COLOR_DARK_BLUE)
    ]
    for i, (title, desc, col_theme) in enumerate(feat_cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(2.1 + row * 2.35)
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        ftb = slide5.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.2), Inches(1.7))
        ftf = ftb.text_frame
        ftf.word_wrap = True
        fp0 = ftf.paragraphs[0]
        fp0.text = title
        fp0.font.size = Pt(16)
        fp0.font.bold = True
        fp0.font.color.rgb = col_theme

        fp1 = ftf.add_paragraph()
        fp1.text = f"\n{desc}"
        fp1.font.size = Pt(13)
        fp1.font.color.rgb = COLOR_BLACK

    # -------------------------------------------------------------
    # SLIDE 6: Market Opportunity & Scalability
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Market & Growth", "Business Model & Expansion Strategy", "Multi-tiered monetization strategy for B2C, B2B, and Enterprise clients.")

    biz_items = [
        ("B2C Candidate Safety Subscriptions", "Free essential checks for students and job seekers; Pro tier ($8/mo) with unlimited batch OCR, continuous domain monitoring & browser extension integration.", COLOR_PRIMARY_BLUE),
        ("B2B Job Board Verification API", "Turnkey REST API integration with platforms like LinkedIn, Unstop, and Internshala to issue authenticated Trust Passports prior to listing.", COLOR_DARK_BLUE),
        ("Enterprise Recruiter Brand Defense", "Brand defense monitoring for Fortune 500 companies to actively detect and takedown spoofed recruiter domains impersonating their hiring teams.", COLOR_PRIMARY_BLUE)
    ]
    for i, (title, desc, col_theme) in enumerate(biz_items):
        left = Inches(0.8 + i * 3.9)
        top = Inches(2.1)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        btb = slide6.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), Inches(3.2), Inches(3.9))
        btf = btb.text_frame
        btf.word_wrap = True
        bp0 = btf.paragraphs[0]
        bp0.text = title
        bp0.font.size = Pt(16)
        bp0.font.bold = True
        bp0.font.color.rgb = col_theme

        bp1 = btf.add_paragraph()
        bp1.text = f"\n{desc}"
        bp1.font.size = Pt(13)
        bp1.font.color.rgb = COLOR_BLACK

    # -------------------------------------------------------------
    # SLIDE 7: Conclusion & Live Demonstration Links
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = COLOR_BG
    bg7.line.fill.background()

    # Blue top bar
    bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    bar7.fill.solid()
    bar7.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
    bar7.line.fill.background()

    card7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.8), Inches(10.933), Inches(5.9))
    card7.fill.solid()
    card7.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    card7.line.color.rgb = COLOR_BLUE_BORDER

    tb7 = slide7.shapes.add_textbox(Inches(1.5), Inches(1.1), Inches(10.333), Inches(5.3))
    tf7 = tb7.text_frame
    tf7.word_wrap = True

    p0 = tf7.paragraphs[0]
    p0.text = "Thank You! Experience HireShield Live"
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_DARK_BLUE
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf7.add_paragraph()
    p1.text = "Protecting the Future of Work with Explainable Cyber Intelligence\n"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY_BLUE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf7.add_paragraph()
    p2.text = "🌐 Live Web Application: https://hireshield-ro6i.onrender.com"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BLACK
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf7.add_paragraph()
    p3.text = "🎥 Demo Video (Google Drive): https://drive.google.com/file/d/1MCPUq-A6OGeyXfKM-25laUY6rRLW4mTh/view?usp=drive_link"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_PRIMARY_BLUE
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf7.add_paragraph()
    p4.text = "📦 GitHub Repository: https://github.com/irksome06/HireShield"
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_MUTED
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf7.add_paragraph()
    p5.text = "\n🔑 Demo Evaluator Credentials: evaluator@hireshield.ai  |  Password: HireShield2026!"
    p5.font.size = Pt(12)
    p5.font.bold = True
    p5.font.color.rgb = COLOR_DARK_BLUE
    p5.alignment = PP_ALIGN.CENTER

    # Save PPTX
    out_path = os.path.join(os.path.dirname(__file__), "HIRESHIELD_PRESENTATION.pptx")
    prs.save(out_path)
    print(f"Presentation generated in White/Black/Blue theme: {out_path}")

if __name__ == "__main__":
    create_presentation()
